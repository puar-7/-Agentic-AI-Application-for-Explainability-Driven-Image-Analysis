# backend/services/file_dispatcher.py
#
# Single responsibility: given a file path, return a list of
# LangChain Document objects.
#
# This module knows nothing about ZIPs, HTTP, or deduplication.
# It is purely concerned with: extension → loader → documents.
#
# Two typed exceptions let callers distinguish between
# "we don't support this type" and "we support it but it broke",
# which map to different response categories (skipped_unsupported
# vs failed) in upload_routes.py.

import os
from typing import List

from langchain_core.documents import Document


# ------------------------------------------------------------------
# Typed exceptions
# ------------------------------------------------------------------

class UnsupportedFileTypeError(Exception):
    """
    Raised when the file extension is not in the supported registry.
    Caller should add the file to skipped_unsupported, not failed.
    """
    pass


class FileLoadError(Exception):
    """
    Raised when the extension is supported but the loader fails.
    Carries a human-readable reason string.
    Caller should add the file to failed with this reason.
    """
    def __init__(self, filename: str, reason: str):
        self.filename = filename
        self.reason = reason
        super().__init__(f"{filename}: {reason}")


# ------------------------------------------------------------------
# Supported extensions registry
# ------------------------------------------------------------------
# Add new extensions here only. The dispatch logic below reads from
# this set for the "is it supported?" check before attempting to load.

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".xlsx", ".pptx"}


# ------------------------------------------------------------------
# Public entry point
# ------------------------------------------------------------------

def load_file(file_path: str) -> List[Document]:
    """
    Load a single file and return a list of LangChain Document objects.

    Args:
        file_path: Absolute or relative path to the file on disk.

    Returns:
        List of Document objects. May be empty if the file has no
        extractable text (e.g. a scanned PDF, an empty spreadsheet).

    Raises:
        UnsupportedFileTypeError: Extension not in SUPPORTED_EXTENSIONS.
        FileLoadError: Extension supported but loading failed.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFileTypeError(
            f"Extension '{ext}' is not supported. "
            f"Supported types: {sorted(SUPPORTED_EXTENSIONS)}"
        )

    # Dispatch to the correct loader
    if ext == ".pdf":
        return _load_pdf(file_path)
    elif ext == ".txt":
        return _load_txt(file_path)
    elif ext == ".docx":
        return _load_docx(file_path)
    elif ext == ".xlsx":
        return _load_xlsx(file_path)
    elif ext == ".pptx":
        return _load_pptx(file_path)

    # Defensive fallback — should never reach here given the check above
    raise UnsupportedFileTypeError(f"No loader registered for '{ext}'")


# ------------------------------------------------------------------
# PDF loader
# ------------------------------------------------------------------

def _load_pdf(file_path: str) -> List[Document]:
    """
    Uses LangChain's PyPDFLoader.
    Returns one Document per page.
    Scanned PDFs (image-only) return Documents with empty page_content —
    this is existing behaviour, not a new problem introduced here.
    """
    try:
        from langchain_community.document_loaders import PyPDFLoader
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        return docs

    except Exception as e:
        raise FileLoadError(
            os.path.basename(file_path),
            f"PDF load failed: {type(e).__name__}: {e}"
        )


# ------------------------------------------------------------------
# TXT loader
# ------------------------------------------------------------------

def _load_txt(file_path: str) -> List[Document]:
    """
    Uses LangChain's TextLoader.

    Encoding strategy:
        1. Try UTF-8 first (most common)
        2. Fall back to latin-1 (reads any byte sequence without error,
           though some characters may be wrong — acceptable for indexing)

    This prevents UnicodeDecodeError from crashing on files produced
    by Windows machines or older research tools.
    """
    try:
        from langchain_community.document_loaders import TextLoader

        try:
            loader = TextLoader(file_path, encoding="utf-8")
            return loader.load()

        except UnicodeDecodeError:
            # Latin-1 can decode any byte sequence — never raises UnicodeDecodeError
            loader = TextLoader(file_path, encoding="latin-1")
            return loader.load()

    except Exception as e:
        raise FileLoadError(
            os.path.basename(file_path),
            f"Text load failed: {type(e).__name__}: {e}"
        )


# ------------------------------------------------------------------
# DOCX loader
# ------------------------------------------------------------------

def _load_docx(file_path: str) -> List[Document]:
    """
    Uses LangChain's Docx2txtLoader (requires docx2txt package).
    Returns the entire document as a single Document object.

    Password-protected .docx files raise an exception which is caught
    and re-raised as FileLoadError so the caller adds it to `failed`.
    """
    try:
        from langchain_community.document_loaders import Docx2txtLoader
        loader = Docx2txtLoader(file_path)
        docs = loader.load()
        return docs

    except ImportError:
        raise FileLoadError(
            os.path.basename(file_path),
            "docx2txt is not installed. Run: pip install docx2txt"
        )

    except Exception as e:
        # Catches password-protected files, corrupt files, etc.
        raise FileLoadError(
            os.path.basename(file_path),
            f"DOCX load failed: {type(e).__name__}: {e}"
        )


# ------------------------------------------------------------------
# XLSX loader (custom — avoids unstructured dependency)
# ------------------------------------------------------------------

def _load_xlsx(file_path: str) -> List[Document]:
    """
    Custom loader using openpyxl directly. Avoids pulling in
    the heavy `unstructured` library.

    Strategy:
        - Each sheet becomes one Document
        - Sheet name is stored in metadata["sheet"]
        - Rows are joined as tab-separated values, one row per line
        - Empty sheets (no data) are skipped silently

    Limitations (accepted):
        - Formula cells show cached values only, not formulas
        - Charts and images are ignored
        - Merged cells show value only in the top-left cell
    """
    try:
        import openpyxl
    except ImportError:
        raise FileLoadError(
            os.path.basename(file_path),
            "openpyxl is not installed. Run: pip install openpyxl"
        )

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        docs = []
        filename = os.path.basename(file_path)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines = []

            for row in ws.iter_rows(values_only=True):
                # Convert each cell to string, replace None with empty string
                row_text = "\t".join(
                    str(cell) if cell is not None else ""
                    for cell in row
                )
                # Skip completely empty rows
                if row_text.strip():
                    lines.append(row_text)

            if not lines:
                # Sheet has no data — skip it silently
                continue

            content = "\n".join(lines)

            docs.append(Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "sheet": sheet_name,
                    "source_type": "document",
                }
            ))

        wb.close()
        return docs

    except Exception as e:
        raise FileLoadError(
            os.path.basename(file_path),
            f"XLSX load failed: {type(e).__name__}: {e}"
        )


# ------------------------------------------------------------------
# PPTX loader (custom — avoids unstructured dependency)
# ------------------------------------------------------------------

def _load_pptx(file_path: str) -> List[Document]:
    """
    Custom loader using python-pptx directly. Avoids pulling in
    the heavy `unstructured` library.

    Strategy:
        - Each slide becomes one Document
        - Slide number is stored in metadata["slide"]
        - All text frames from all shapes on the slide are extracted
        - Slides with no extractable text are skipped silently

    Limitations (accepted):
        - Image-only slides produce no text (common in design-heavy decks)
        - SmartArt text may not be fully extracted
        - Slide notes are not included (can be added later if needed)
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise FileLoadError(
            os.path.basename(file_path),
            "python-pptx is not installed. Run: pip install python-pptx"
        )

    try:
        prs = Presentation(file_path)
        docs = []
        filename = os.path.basename(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []

            for shape in slide.shapes:
                # Only process shapes that have a text frame
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(
                        run.text for run in paragraph.runs
                    ).strip()

                    if para_text:
                        slide_texts.append(para_text)

            if not slide_texts:
                # Slide has no extractable text — skip silently
                continue

            content = "\n".join(slide_texts)

            docs.append(Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "slide": slide_num,
                    "source_type": "document",
                }
            ))

        return docs

    except Exception as e:
        raise FileLoadError(
            os.path.basename(file_path),
            f"PPTX load failed: {type(e).__name__}: {e}"
        )